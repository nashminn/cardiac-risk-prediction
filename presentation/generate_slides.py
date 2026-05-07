#!/usr/bin/env python3
"""Generate cardiac risk prediction PPTX presentation (Google Slides compatible)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE    = "/home/nashmin/Documents/exp/cardiac-risk-prediction"
IMG_DIR = os.path.join(BASE, "outputs")
OUT_DIR = os.path.join(BASE, "presentation")
OUT     = os.path.join(OUT_DIR, "cardiac_risk_prediction.pptx")
os.makedirs(OUT_DIR, exist_ok=True)

# Slide dimensions: 16:9 widescreen
SW = Inches(13.333)
SH = Inches(7.5)

# Color palette (light, projector-friendly)
NAVY  = RGBColor(0x1A, 0x3A, 0x5C)
BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
LBLUE = RGBColor(0xD5, 0xE5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1C, 0x1C, 0x2E)
GRAY  = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0xF2, 0xF2, 0xF2)
RED   = RGBColor(0xBE, 0x3A, 0x2A)
GREEN = RGBColor(0x1E, 0x7A, 0x3E)
AMBER = RGBColor(0xD4, 0x6B, 0x08)

HDR_H = Inches(1.15)
PAD   = Inches(0.55)
CT    = HDR_H + Inches(0.15)   # content top (1.3 in)
N     = 12                     # numbered content slides

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH


# ── helpers ───────────────────────────────────────────────────────────

def new_slide():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    return sl


def add_rect(sl, l, t, w, h, fill, border=False):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = BLUE
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def add_text(sl, l, t, w, h, txt, sz, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tx.word_wrap = True
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tx


def add_bullets(sl, l, t, w, h, items, sz=15, base_color=DARK):
    """items: str  or  (text, level, bold, color)  — trailing fields optional."""
    tx = sl.shapes.add_textbox(l, t, w, h)
    tx.word_wrap = True
    tf = tx.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, str):
            txt, lvl, bold, col = it, 0, False, base_color
        else:
            parts = list(it) + [None] * 4
            txt  = parts[0]
            lvl  = parts[1] if parts[1] is not None else 0
            bold = parts[2] if parts[2] is not None else False
            col  = parts[3] if parts[3] is not None else base_color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        prefix = ("    " * lvl) + ("● " if lvl == 0 else "– ")
        r = p.add_run()
        r.text = prefix + txt
        r.font.size = Pt(sz)
        r.font.bold = bool(bold)
        r.font.color.rgb = col if col else base_color


def add_header(sl, title, num=None):
    add_rect(sl, 0, 0, SW, HDR_H, LBLUE)
    add_rect(sl, 0, HDR_H - Inches(0.05), SW, Inches(0.05), BLUE)
    add_text(sl, PAD, Inches(0.22), SW - Inches(2), Inches(0.85),
             title, 22, bold=True, color=NAVY)
    if num is not None:
        add_text(sl, SW - Inches(1.65), SH - Inches(0.44), Inches(1.5), Inches(0.4),
                 f"{num}/{N}", 12, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image(sl, fname, l, t, w, h=None):
    path = os.path.join(IMG_DIR, fname)
    if h is not None:
        return sl.shapes.add_picture(path, l, t, w, h)
    return sl.shapes.add_picture(path, l, t, w)


# ── SLIDE 1 — TITLE ──────────────────────────────────────────────────
sl = new_slide()

add_rect(sl, 0, 0, SW, Inches(3.0), NAVY)
add_rect(sl, 0, Inches(3.0), SW, Inches(0.12), BLUE)

add_text(sl, 0, Inches(0.32), SW, Inches(0.52),
         "CSE710: Advanced Artificial Intelligence",
         14, color=LBLUE, align=PP_ALIGN.CENTER)

add_text(sl, Inches(0.8), Inches(0.82), SW - Inches(1.6), Inches(1.35),
         "Interpretable Cardiac Risk Prediction",
         36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, Inches(0.8), Inches(2.08), SW - Inches(1.6), Inches(0.75),
         "A Deep Learning Approach with SHAP-based Clinical Explanations",
         18, italic=True, color=LBLUE, align=PP_ALIGN.CENTER)

# Author card
add_rect(sl, Inches(3.4), Inches(3.35), Inches(6.5), Inches(3.1), LGRAY, border=True)
add_text(sl, Inches(3.4), Inches(3.52), Inches(6.5), Inches(0.52),
         "Project Team", 14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

for i, name in enumerate(["[Person 1 Name]", "[Person 2 Name]", "[Person 3 Name]"]):
    add_text(sl, Inches(3.4), Inches(4.1) + Inches(0.6) * i,
             Inches(6.5), Inches(0.52),
             name, 18, color=DARK, align=PP_ALIGN.CENTER)

add_text(sl, 0, SH - Inches(0.62), SW, Inches(0.48),
         "May 2026", 13, color=GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — OUTLINE ────────────────────────────────────────────────
sl = new_slide()
add_header(sl, "Presentation Outline")

outline = [
    "Problem Statement & Motivation",
    "Dataset Overview",
    "Data Preprocessing",
    "Exploratory Data Analysis",
    "MLP Architecture",
    "Training & Optimization",
    "Model Evaluation",
    "Explainable AI — SHAP",
    "Key Findings & Limitations",
    "References",
]

col1, col2 = outline[:5], outline[5:]
cy = CT + Inches(0.1)
for item in col1:
    add_text(sl, PAD, cy, Inches(5.8), Inches(0.55), f"● {item}", 17, color=DARK)
    cy += Inches(0.62)

cy = CT + Inches(0.1)
for item in col2:
    add_text(sl, Inches(7.2), cy, Inches(5.7), Inches(0.55), f"● {item}", 17, color=DARK)
    cy += Inches(0.62)


# ── SLIDE 3 — PROBLEM STATEMENT (1/12) ──────────────────────────────
sl = new_slide()
add_header(sl, "Problem Statement & Motivation", num=1)

add_text(sl, PAD, CT, Inches(7.5), Inches(0.48),
         "Why Cardiac Risk Prediction Matters", 17, bold=True, color=NAVY)

add_bullets(sl, PAD, CT + Inches(0.52), Inches(7.5), Inches(1.8), [
    "Cardiovascular disease is the #1 cause of death globally (WHO, 2023)",
    "Early detection dramatically reduces mortality and healthcare costs",
    "Routine clinical data (ECG, cholesterol, blood pressure) already holds predictive power",
], sz=15)

add_text(sl, PAD, CT + Inches(2.42), Inches(7.5), Inches(0.48),
         "The \"Black Box\" Problem in Clinical AI", 17, bold=True, color=NAVY)

add_bullets(sl, PAD, CT + Inches(2.94), Inches(7.5), Inches(1.9), [
    "Deep learning is powerful but opaque — clinicians cannot trust what they cannot understand",
    "Regulatory frameworks (FDA, EU AI Act) increasingly require explainability for medical AI",
    "Without interpretability, accurate predictions may still be rejected in clinical practice",
], sz=15)

# Right panel
add_rect(sl, Inches(8.55), CT, Inches(4.4), Inches(5.5), LGRAY, border=True)
add_text(sl, Inches(8.65), CT + Inches(0.12), Inches(4.2), Inches(0.48),
         "Our Solution", 15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_bullets(sl, Inches(8.7), CT + Inches(0.65), Inches(3.9), Inches(2.8), [
    ("MLP deep learning model", 0, True, BLUE),
    ("SHAP for per-prediction explanations", 0, True, BLUE),
    "UCI Heart Disease dataset",
    "302 patients · 13 clinical features",
    "Binary target: Cardiac Risk / No Risk",
], sz=14)

add_text(sl, Inches(8.7), CT + Inches(3.7), Inches(3.9), Inches(1.6),
         "Goal: an accurate AND interpretable model that clinicians can act on.",
         14, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── SLIDE 4 — DATASET OVERVIEW (2/12) ────────────────────────────────
sl = new_slide()
add_header(sl, "Dataset Overview — UCI Heart Disease (Cleveland)", num=2)

add_text(sl, PAD, CT, Inches(6.5), Inches(0.46),
         "Dataset at a Glance", 16, bold=True, color=NAVY)

add_bullets(sl, PAD, CT + Inches(0.5), Inches(6.3), Inches(1.9), [
    "Source: UCI Machine Learning Repository — Cleveland subset",
    "Curated on Kaggle by mragpavank",
    "303 records → 302 after removing 1 duplicate",
    "Classes: 164 Cardiac Risk (54%)  ·  138 No Risk (46%)",
], sz=14)

# Feature mini-table
features = [
    ("Age",            "Age in years"),
    ("ChestPainType",  "0 = asymptomatic … 3 = typical angina"),
    ("RestingBP",      "Resting blood pressure (mm Hg)"),
    ("Cholesterol",    "Serum cholesterol (mg/dl)"),
    ("MaxHR",          "Maximum heart rate achieved"),
    ("Oldpeak",        "ST depression induced by exercise"),
    ("Thalassemia",    "Thalassemia type (0–3)"),
]

tbl_top = CT + Inches(2.5)
add_text(sl, PAD, tbl_top - Inches(0.38), Inches(6.5), Inches(0.38),
         "Selected Features (13 total)", 14, bold=True, color=NAVY)

rh = Inches(0.44)
cw1, cw2 = Inches(2.05), Inches(4.1)
for idx, (feat, desc) in enumerate(features):
    bg = LGRAY if idx % 2 == 0 else WHITE
    add_rect(sl, PAD, tbl_top + idx * rh, cw1 + cw2, rh, bg, border=True)
    add_text(sl, PAD + Inches(0.06), tbl_top + idx * rh + Inches(0.07),
             cw1 - Inches(0.1), rh, feat, 12, bold=True, color=NAVY)
    add_text(sl, PAD + cw1 + Inches(0.06), tbl_top + idx * rh + Inches(0.07),
             cw2 - Inches(0.1), rh, desc, 12, color=DARK)
# table bottom: 1.3 + 2.5 + 7*0.44 = 1.3+2.5+3.08 = 6.88 ✓

# Right: class distribution
add_text(sl, Inches(7.05), CT, Inches(5.9), Inches(0.46),
         "Class Distribution", 15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_image(sl, "class_distribution.png", Inches(7.35), CT + Inches(0.5), Inches(5.4))


# ── SLIDE 5 — PREPROCESSING (3/12) ──────────────────────────────────
sl = new_slide()
add_header(sl, "Data Preprocessing Pipeline", num=3)

steps = [
    ("Step 1 — Missing Values",
     "No missing values found across all 303 records and 14 columns."),
    ("Step 2 — Deduplication",
     "1 exact duplicate row removed → 302 clean samples retained."),
    ("Step 3 — Feature / Target Split",
     "X = 13 clinical features.   y = binary target (1 = Cardiac Risk, 0 = No Risk)."),
    ("Step 4 — Train-Test Split",
     "80% Training (241 samples)  ·  20% Testing (61 samples)  ·  Stratified by class label."),
    ("Step 5 — Standardization",
     "StandardScaler fitted on training data only; applied to both train and test sets."),
]

cy = CT + Inches(0.1)
for title, desc in steps:
    add_rect(sl, PAD, cy, SW - PAD * 2, Inches(0.98), LGRAY, border=True)
    add_text(sl, PAD + Inches(0.15), cy + Inches(0.07),
             Inches(3.8), Inches(0.44), title, 14, bold=True, color=BLUE)
    add_text(sl, PAD + Inches(0.15), cy + Inches(0.48),
             SW - PAD * 2 - Inches(0.3), Inches(0.42), desc, 13, color=DARK)
    cy += Inches(1.07)
# bottom: 1.3 + 0.1 + 5*1.07 = 6.75 ✓

add_text(sl, PAD, cy + Inches(0.08), SW - PAD * 2, Inches(0.48),
         "Stratification ensures class proportions are preserved; "
         "fitting the scaler only on training data prevents data leakage.",
         13, italic=True, color=GRAY)


# ── SLIDE 6 — EDA (4/12) ────────────────────────────────────────────
sl = new_slide()
add_header(sl, "Exploratory Data Analysis", num=4)

add_text(sl, PAD, CT, Inches(6.4), Inches(0.43),
         "Correlation Heatmap", 15, bold=True, color=NAVY)
# heatmap: figsize (12,9) → at 6.4" wide: 6.4*9/12 = 4.8" tall → top 1.73, bottom 6.53 ✓
add_image(sl, "correlation_heatmap.png", PAD, CT + Inches(0.45), Inches(6.4))

add_text(sl, Inches(7.15), CT, Inches(5.8), Inches(0.43),
         "Class Distribution by Feature", 15, bold=True, color=NAVY)
# class_dist: figsize (5,4) → at 5.5" wide: 5.5*4/5 = 4.4" tall → top 1.73, bottom 6.13 ✓
add_image(sl, "class_distribution.png", Inches(7.55), CT + Inches(0.45), Inches(5.3))

add_text(sl, PAD, SH - Inches(0.55), SW - PAD * 2, Inches(0.42),
         "Key: MaxHR negatively correlated with risk · Oldpeak strongly positive · "
         "Thalassemia and ChestPainType show high discriminative power.",
         12, italic=True, color=GRAY)


# ── SLIDE 7 — MLP ARCHITECTURE (5/12) ───────────────────────────────
sl = new_slide()
add_header(sl, "MLP Architecture", num=5)

# Left: architecture layers
add_text(sl, PAD, CT, Inches(5.8), Inches(0.46),
         "Network Design", 16, bold=True, color=NAVY)

arch_layers = [
    ("Input Layer",       "13 features (standardized)",           BLUE),
    ("Dense 64 + ReLU",   "First hidden layer + Dropout(0.3)",    BLUE),
    ("Dense 32 + ReLU",   "Second hidden layer + Dropout(0.3)",   BLUE),
    ("Dense 16 + ReLU",   "Third hidden layer",                   BLUE),
    ("Output — Sigmoid",  "P(Cardiac Risk)  ∈  [0, 1]",          GREEN),
]

cy = CT + Inches(0.5)
box_h = Inches(0.82)
for layer, desc, col in arch_layers:
    add_rect(sl, PAD, cy, Inches(5.8), box_h, LGRAY, border=True)
    add_text(sl, PAD + Inches(0.12), cy + Inches(0.06),
             Inches(2.3), Inches(0.42), layer, 13, bold=True, color=col)
    add_text(sl, PAD + Inches(0.12), cy + Inches(0.44),
             Inches(5.5), Inches(0.34), desc, 13, color=DARK)
    cy += Inches(0.9)

add_text(sl, PAD, cy + Inches(0.1), Inches(5.8), Inches(0.46),
         "Total Parameters: 3,521  (13.75 KB)", 14, bold=True, color=GREEN)

# Right: training config
add_text(sl, Inches(6.8), CT, Inches(6.2), Inches(0.46),
         "Training Configuration", 16, bold=True, color=NAVY)

add_bullets(sl, Inches(6.8), CT + Inches(0.52), Inches(6.2), Inches(3.2), [
    "Optimizer: Adam (adaptive learning rate)",
    "Loss Function: Binary Cross-Entropy",
    "Batch Size: 16  ·  Max Epochs: 200",
    "Validation Split: 15% of training set",
    "Early Stopping: patience = 15 on val_loss",
    "Best weights automatically restored on stop",
], sz=15)

add_rect(sl, Inches(6.8), CT + Inches(3.85), Inches(6.2), Inches(1.9),
         RGBColor(0xFF, 0xF5, 0xD6), border=True)
add_text(sl, Inches(6.95), CT + Inches(3.97), Inches(5.8), Inches(1.65),
         "Why Dropout? With only 302 samples, the model easily memorises the training set. "
         "Randomly zeroing 30% of neurons each step forces the network to learn redundant "
         "representations — acting as an ensemble of many sub-networks.",
         13, italic=True, color=DARK)


# ── SLIDE 8 — TRAINING (6/12) ────────────────────────────────────────
sl = new_slide()
add_header(sl, "Training & Optimization", num=6)

add_text(sl, PAD, CT, Inches(7.1), Inches(0.44),
         "Training History", 15, bold=True, color=NAVY)
# training_history: figsize (12,4) → at 7.0" wide: 7.0*4/12 = 2.33" tall → bottom 4.07 ✓
add_image(sl, "training_history.png", PAD, CT + Inches(0.46), Inches(7.1))

add_text(sl, Inches(8.0), CT, Inches(5.0), Inches(0.44),
         "Observations", 15, bold=True, color=NAVY)

add_bullets(sl, Inches(8.0), CT + Inches(0.5), Inches(5.0), Inches(3.2), [
    "Stopped at epoch ~33 via early stopping",
    "val_loss plateaued after epoch ~14",
    "No significant overfitting observed",
    "Best weights from epoch ~14 restored",
    ("Small train/val accuracy gap expected at n=302", 1, False, GRAY),
], sz=14)

add_rect(sl, Inches(8.0), CT + Inches(3.75), Inches(5.0), Inches(2.0),
         LGRAY, border=True)
add_text(sl, Inches(8.15), CT + Inches(3.88), Inches(4.7), Inches(1.78),
         "Early stopping is essential here — without it, 200 epochs on 241 samples "
         "would cause severe overfitting. The model converged in ~33 epochs, "
         "saving compute and preserving generalisation.",
         13, italic=True, color=DARK)


# ── SLIDE 9 — EVALUATION METRICS (7/12) ─────────────────────────────
sl = new_slide()
add_header(sl, "Model Evaluation — Performance Metrics", num=7)

metrics = [
    ("Accuracy",  "80.33%", "Overall correct predictions",        BLUE),
    ("Precision", "76.92%", "Of predicted Risk: truly Risk",      BLUE),
    ("Recall",    "90.91%", "Of true Risk cases: caught",         GREEN),
    ("F1-Score",  "83.33%", "Harmonic mean of Prec & Recall",     GREEN),
    ("ROC-AUC",   "86.90%", "Area under the ROC curve",           GREEN),
]

bw = Inches(2.27)
gap = Inches(0.14)
total_w = len(metrics) * bw + (len(metrics) - 1) * gap
sx = (SW - total_w) / 2

for i, (name, val, desc, col) in enumerate(metrics):
    lx = sx + i * (bw + gap)
    add_rect(sl, lx, CT + Inches(0.25), bw, Inches(2.3), LGRAY, border=True)
    add_text(sl, lx, CT + Inches(0.35), bw, Inches(0.5),
             name, 14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(sl, lx, CT + Inches(0.85), bw, Inches(0.85),
             val, 30, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, lx + Inches(0.06), CT + Inches(1.72), bw - Inches(0.12), Inches(0.76),
             desc, 11, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(sl, PAD, CT + Inches(2.8), SW - PAD * 2, Inches(1.75),
         RGBColor(0xE8, 0xF5, 0xEB), border=True)
add_text(sl, PAD + Inches(0.2), CT + Inches(2.95), SW - PAD * 2 - Inches(0.4), Inches(1.5),
         "Clinical interpretation: Recall (90.91%) is the most critical metric in cardiac "
         "screening. A false negative — missing a real cardiac risk patient — is far more "
         "dangerous than a false alarm. The model correctly catches 30 of 33 true risk "
         "cases in the test set (only 3 missed).",
         14, color=DARK)

add_text(sl, PAD, CT + Inches(4.7), SW - PAD * 2, Inches(0.46),
         "Confusion Matrix Detail:   TP = 30   ·   TN = 19   ·   FP = 9   ·   FN = 3",
         14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── SLIDE 10 — EVALUATION VISUALS (8/12) ────────────────────────────
sl = new_slide()
add_header(sl, "Model Evaluation — Confusion Matrix & ROC Curve", num=8)

add_text(sl, PAD, CT, Inches(5.8), Inches(0.43),
         "Confusion Matrix", 15, bold=True, color=NAVY)
# confusion_matrix: figsize (5,4) → at 5.5" wide: 5.5*4/5 = 4.4" tall → bottom 6.13 ✓
add_image(sl, "confusion_matrix.png", PAD, CT + Inches(0.45), Inches(5.5))

add_text(sl, Inches(6.7), CT, Inches(6.2), Inches(0.43),
         "ROC Curve  (AUC = 0.8690)", 15, bold=True, color=NAVY)
# roc_curve: figsize (6,5) → at 6.2" wide: 6.2*5/6 = 5.17" tall → bottom 6.95 — tight
add_image(sl, "roc_curve.png", Inches(6.75), CT + Inches(0.45), Inches(6.0))


# ── SLIDE 11 — SHAP OVERVIEW (9/12) ─────────────────────────────────
sl = new_slide()
add_header(sl, "Explainable AI — What is SHAP?", num=9)

add_text(sl, PAD, CT, SW - PAD * 2, Inches(0.48),
         "SHAP  —  SHapley Additive exPlanations", 18, bold=True, color=NAVY)

add_text(sl, PAD, CT + Inches(0.52), SW - PAD * 2, Inches(0.42),
         "Rooted in cooperative game theory (Shapley, 1953)  ·  "
         "Adapted for ML by Lundberg & Lee (NeurIPS 2017)",
         14, italic=True, color=GRAY)

# Three axioms
props = [
    ("Local Accuracy",
     "For any prediction, the SHAP values sum exactly to the model output minus a baseline. "
     "Explanation and model always agree."),
    ("Missingness",
     "Features that are absent from the input have zero SHAP contribution — no phantom effects."),
    ("Consistency",
     "If a feature's contribution increases across all coalition subsets, "
     "its SHAP value can only increase — a monotonicity guarantee."),
]
cy = CT + Inches(1.05)
for prop, desc in props:
    add_rect(sl, PAD, cy, SW - PAD * 2, Inches(1.02), LGRAY, border=True)
    add_text(sl, PAD + Inches(0.15), cy + Inches(0.08),
             Inches(3.5), Inches(0.46), prop, 14, bold=True, color=BLUE)
    add_text(sl, PAD + Inches(0.15), cy + Inches(0.5),
             SW - PAD * 2 - Inches(0.3), Inches(0.44), desc, 13, color=DARK)
    cy += Inches(1.1)
# bottom: 1.3+1.05+3*1.1 = 5.65 ✓

add_text(sl, PAD, cy + Inches(0.15), SW - PAD * 2, Inches(0.46),
         "Why KernelExplainer?", 16, bold=True, color=NAVY)

add_bullets(sl, PAD, cy + Inches(0.65), SW - PAD * 2, Inches(1.0), [
    "Model-agnostic — works with any black-box, including TensorFlow/Keras MLP",
    "Samples feature perturbations to estimate Shapley values without model-specific access",
    "Setup: 100 background samples  ·  50 test samples  ·  nsamples=100 per explanation",
], sz=14)


# ── SLIDE 12 — SHAP GLOBAL (10/12) ──────────────────────────────────
sl = new_slide()
add_header(sl, "SHAP — Global Feature Importance", num=10)

add_text(sl, PAD, CT, Inches(6.2), Inches(0.43),
         "Mean |SHAP| — Feature Importance Bar", 15, bold=True, color=NAVY)
# shap_feature_importance: default matplotlib ~6.4x4.8 → at 6.1" wide: 4.57" tall → bottom 6.3 ✓
add_image(sl, "shap_feature_importance.png", PAD, CT + Inches(0.46), Inches(6.1))

add_text(sl, Inches(6.95), CT, Inches(6.0), Inches(0.43),
         "Beeswarm — Direction & Magnitude", 15, bold=True, color=NAVY)
# shap_summary_plot: default ~6.4x4.8 → at 5.9" wide: 4.42" tall → bottom 6.18 ✓
add_image(sl, "shap_summary_plot.png", Inches(6.95), CT + Inches(0.46), Inches(5.9))

add_text(sl, PAD, SH - Inches(0.52), SW - PAD * 2, Inches(0.42),
         "Top global drivers: Thalassemia  ·  ChestPainType  ·  MaxHR  ·  Oldpeak  ·  MajorVessels",
         13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 13 — SHAP LOCAL (11/12) ───────────────────────────────────
sl = new_slide()
add_header(sl, "SHAP — Local Explanation (Patient-Level)", num=11)

add_text(sl, PAD, CT, Inches(6.6), Inches(0.43),
         "Waterfall Plot — Patient 1", 15, bold=True, color=NAVY)
# shap_local_explanation: default ~6.4x4.8 → at 6.5" wide: 4.875" → bottom 6.65 — OK
add_image(sl, "shap_local_explanation.png", PAD, CT + Inches(0.46), Inches(6.5))

add_text(sl, Inches(7.4), CT, Inches(5.55), Inches(0.46),
         "Reading the Waterfall", 16, bold=True, color=NAVY)

add_bullets(sl, Inches(7.4), CT + Inches(0.5), Inches(5.5), Inches(3.3), [
    "Patient 1: Actual = No Risk · Predicted = No Risk  (P = 0.094)  ✓",
    "MaxHR drives prediction strongly toward No Risk (large negative SHAP)",
    "Thalassemia also reduces predicted risk for this patient",
    "ChestPainType adds minor push toward risk",
    ("Each bar moves the output from the baseline expected value (~0.54)", 1, False, GRAY),
], sz=14)

add_rect(sl, Inches(7.4), CT + Inches(3.95), Inches(5.55), Inches(1.75),
         LGRAY, border=True)
add_text(sl, Inches(7.55), CT + Inches(4.08), Inches(5.2), Inches(1.55),
         "Unlike global rankings, SHAP local explanations let a clinician see exactly "
         "why the model flagged a specific patient — enabling audit, override, "
         "and trust-building in practice.",
         13, italic=True, color=DARK)


# ── SLIDE 14 — KEY FINDINGS (12/12) ─────────────────────────────────
sl = new_slide()
add_header(sl, "Key Findings & Limitations", num=12)

add_text(sl, PAD, CT, Inches(6.2), Inches(0.46),
         "Key Findings", 16, bold=True, color=NAVY)

add_bullets(sl, PAD, CT + Inches(0.5), Inches(6.2), Inches(3.8), [
    ("MLP achieves 80.33% accuracy and 90.91% recall on 302 samples", 0, True, DARK),
    "High recall minimises false negatives — the most critical error in cardiac screening",
    "Thalassemia and ChestPainType are the strongest global predictors",
    "MaxHR is protective — higher max heart rate lowers predicted risk",
    "Local SHAP explanations make individual predictions auditable",
], sz=14)

add_text(sl, Inches(7.2), CT, Inches(5.7), Inches(0.46),
         "Limitations", 16, bold=True, color=RED)

add_bullets(sl, Inches(7.2), CT + Inches(0.5), Inches(5.7), Inches(2.15), [
    "Small dataset (302 samples) limits generalisation",
    "Cleveland subset only — may not reflect diverse populations",
    "KernelExplainer is slow for large test sets (sampling-based)",
    "Binary target collapses 5 original severity levels into 2",
], sz=14)

add_text(sl, Inches(7.2), CT + Inches(2.8), Inches(5.7), Inches(0.46),
         "Future Work", 16, bold=True, color=GREEN)

add_bullets(sl, Inches(7.2), CT + Inches(3.3), Inches(5.7), Inches(2.4), [
    "Larger, multi-centre datasets for better generalisation",
    "TreeSHAP or GradientSHAP for faster local explanations",
    "Extend to multi-class severity prediction",
    "Real-time clinical decision-support integration",
], sz=14)


# ── SLIDE 15 — REFERENCES ────────────────────────────────────────────
sl = new_slide()
add_header(sl, "References")

refs = [
    "[1]  Detrano, R., et al. (1989). International application of a new probability algorithm for "
    "the diagnosis of coronary artery disease. American Journal of Cardiology, 64(5), 304–310.",

    "[2]  Dua, D. & Graff, C. (2019). UCI Machine Learning Repository. "
    "University of California, Irvine.",

    "[3]  mragpavank. (2020). Heart Disease UCI. Kaggle. "
    "kaggle.com/code/mragpavank/heart-disease-uci",

    "[4]  Lundberg, S. M. & Lee, S.-I. (2017). A unified approach to interpreting model predictions. "
    "Advances in Neural Information Processing Systems (NeurIPS), 30.",

    "[5]  Shapley, L. S. (1953). A value for n-person games. "
    "Contributions to the Theory of Games, 2, 307–317.",

    "[6]  Abadi, M., et al. (2015). TensorFlow: Large-scale machine learning on heterogeneous "
    "distributed systems. arXiv:1603.04467.",

    "[7]  Chollet, F., et al. (2015). Keras. GitHub: github.com/keras-team/keras",

    "[8]  World Health Organization. (2023). Cardiovascular diseases fact sheet. WHO.",
]

cy = CT + Inches(0.05)
for ref in refs:
    add_text(sl, PAD, cy, SW - PAD * 2, Inches(0.62), ref, 12, color=DARK)
    cy += Inches(0.66)
# 8 refs * 0.66 = 5.28 → bottom: 1.3+0.05+5.28 = 6.63 ✓


# ── SLIDE 16 — THANK YOU ─────────────────────────────────────────────
sl = new_slide()

add_rect(sl, 0, 0, SW, SH, NAVY)
add_rect(sl, 0, Inches(3.15), SW, Inches(0.12), BLUE)

add_text(sl, 0, Inches(0.9), SW, Inches(1.5),
         "Thank You", 54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, 0, Inches(2.55), SW, Inches(0.85),
         "Any Questions?", 30, italic=True, color=LBLUE, align=PP_ALIGN.CENTER)

add_text(sl, 0, Inches(3.9), SW, Inches(0.6),
         "CSE710: Advanced Artificial Intelligence", 16, color=LBLUE, align=PP_ALIGN.CENTER)

add_text(sl, 0, Inches(4.55), SW, Inches(0.6),
         "[Person 1 Name]   ·   [Person 2 Name]   ·   [Person 3 Name]",
         16, color=LBLUE, align=PP_ALIGN.CENTER)

add_text(sl, 0, Inches(5.45), SW, Inches(0.55),
         "Dataset: UCI Heart Disease (Cleveland)  ·  Model: MLP  ·  XAI: SHAP",
         13, color=GRAY, align=PP_ALIGN.CENTER)


# ── SAVE ─────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✓ Saved {prs.slides.__len__()} slides → {OUT}")
